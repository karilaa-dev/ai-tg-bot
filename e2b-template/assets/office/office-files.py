#!/opt/office/python/bin/python
"""Read-only Office package checks and actual-file rendering. Never rewrite the source."""
import argparse
import fcntl
import hashlib
import json
import os
from pathlib import Path
import posixpath
import re
import resource
import shutil
import subprocess
import sys
import time
from urllib.parse import quote, unquote, urlsplit, urlunsplit
import zipfile
from lxml import etree

MAX_EXPANDED = 200 * 1024 * 1024
OFFICE_RELATIONSHIPS = (
    'http://schemas.openxmlformats.org/officeDocument/2006/relationships/',
    'http://purl.oclc.org/ooxml/officeDocument/relationships/',
)
HYPERLINK_ELEMENTS = {
    '{' + namespace + '}' + tag
    for namespace, tags in [
        ('http://schemas.openxmlformats.org/wordprocessingml/2006/main', ['hyperlink']),
        ('http://purl.oclc.org/ooxml/wordprocessingml/main', ['hyperlink']),
        ('http://schemas.openxmlformats.org/spreadsheetml/2006/main', ['hyperlink']),
        ('http://purl.oclc.org/ooxml/spreadsheetml/main', ['hyperlink']),
        ('http://schemas.openxmlformats.org/drawingml/2006/main', ['hlinkClick', 'hlinkMouseOver']),
        ('http://purl.oclc.org/ooxml/drawingml/main', ['hlinkClick', 'hlinkMouseOver']),
    ] for tag in tags
}


def relationship_target(base, target):
    uri = urlsplit(target)
    if uri.scheme or uri.netloc:
        raise ValueError('Internal relationship contains an external URI')
    path = unquote(uri.path)
    return posixpath.normpath(path.lstrip('/') if path.startswith('/') else posixpath.join(base, path))


def run(args, timeout=120):
    result = subprocess.run(args, capture_output=True, text=True, timeout=timeout)
    if result.returncode:
        raise ValueError((result.stderr or result.stdout or f'{args[0]} exited {result.returncode}')[-4000:])
    return result.stdout


def package_checks(source):
    with zipfile.ZipFile(source) as archive:
        infos = archive.infolist()
        names = set(archive.namelist())
        if len(names) != len(infos) or len(names) > 10000:
            raise ValueError('Duplicate or excessive package entries')
        if sum(i.file_size for i in infos) > MAX_EXPANDED:
            raise ValueError('Expanded package exceeds 200 MiB')
        if archive.testzip():
            raise ValueError('Package CRC check failed')
        main = {'.docx': 'word/document.xml', '.pptx': 'ppt/presentation.xml', '.xlsx': 'xl/workbook.xml'}[source.suffix.lower()]
        for name in ['[Content_Types].xml', '_rels/.rels', main]:
            if name not in names:
                raise ValueError(f'Missing required part: {name}')
        xml = {}
        for name in names:
            if name.endswith(('.xml', '.rels')):
                data = archive.read(name)
                parser = etree.XMLParser(resolve_entities=False, no_network=True, load_dtd=False)
                root = etree.fromstring(data, parser)
                if root.getroottree().docinfo.doctype:
                    raise ValueError(f'DTD is unsupported: {name}')
                xml[name] = root
        relationships = {}
        external_ids = {}
        for name, root in xml.items():
            if not name.endswith('.rels'):
                continue
            base = posixpath.dirname(posixpath.dirname(name)) if name != '_rels/.rels' else ''
            ids = set()
            relationships[name] = ids
            external_ids[name] = set()
            for rel in root:
                rid = rel.get('Id')
                if not rid or rid in ids:
                    raise ValueError(f'Missing or duplicate relationship Id in {name}')
                ids.add(rid)
                if rel.get('TargetMode') == 'External':
                    external_ids[name].add(rid)
                    # Clickable web/email links are inert during rendering. Linked media,
                    # templates, workbooks, and unknown relationship types must not load.
                    if (rel.get('Type') not in [prefix + 'hyperlink' for prefix in OFFICE_RELATIONSHIPS]
                            or urlsplit(rel.get('Target', '')).scheme.lower() not in ('https', 'http', 'mailto')):
                        raise ValueError(f'External resource is unsupported: {name} ({rel.get("Type", "unknown")}). Embed the resource before validation.')
                    continue
                target = rel.get('Target', '')
                resolved = relationship_target(base, target)
                if resolved not in names:
                    raise ValueError(f'Missing relationship target: {name} -> {target}')
        office_rel = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
        strict_rel = 'http://purl.oclc.org/ooxml/officeDocument/relationships'
        for name, root in xml.items():
            if name.endswith('.rels'):
                continue
            relpath = posixpath.join(posixpath.dirname(name), '_rels', posixpath.basename(name) + '.rels')
            for element in root.iter():
                for attribute, value in element.attrib.items():
                    if attribute.startswith(('{' + office_rel + '}', '{' + strict_rel + '}')):
                        if value not in relationships.get(relpath, set()):
                            raise ValueError(f'Missing relationship Id: {name} -> {value}')
                        if value in external_ids.get(relpath, set()) and (element.tag not in HYPERLINK_ELEMENTS or etree.QName(attribute).localname != 'id'):
                            raise ValueError(f'External resource referenced outside a hyperlink: {name} -> {value}')
        if not any(rel.get('Type') in [prefix + 'officeDocument' for prefix in OFFICE_RELATIONSHIPS]
                   and rel.get('TargetMode') != 'External'
                   and relationship_target('', rel.get('Target', '')) == main for rel in xml['_rels/.rels']):
            raise ValueError('Missing main document relationship')
        ns = {'ct': 'http://schemas.openxmlformats.org/package/2006/content-types'}
        ct = xml['[Content_Types].xml']
        overrides = {unquote(e.get('PartName', '')).lstrip('/'): e.get('ContentType') for e in ct.findall('ct:Override', ns)}
        defaults = {e.get('Extension'): e.get('ContentType') for e in ct.findall('ct:Default', ns)}
        for name in names:
            if name.endswith('/') or name == '[Content_Types].xml':
                continue
            if name not in overrides and name.rsplit('.', 1)[-1] not in defaults:
                raise ValueError(f'Missing content type for {name}')
        for name in overrides:
            if name not in names:
                raise ValueError(f'Content type refers to missing part: {name}')
        return len(names)


def convert(source, out, profile, extension):
    run(['libreoffice', '-env:UserInstallation=' + profile.as_uri(), '--headless', '--nologo', '--nodefault', '--norestore', '--convert-to', extension, '--outdir', str(out), str(source)])
    target = out / (source.stem + '.' + extension.split(':')[0])
    if not target.is_file() or not target.stat().st_size:
        raise ValueError(f'LibreOffice did not produce {target.name}')
    return target


def conversion_copy(source, out):
    """Canonicalize equivalent internal URIs for LibreOffice, preserving all other parts."""
    changed = {}
    with zipfile.ZipFile(source) as archive:
        for name in archive.namelist():
            if not name.endswith('.rels'):
                continue
            root = etree.fromstring(archive.read(name), etree.XMLParser(resolve_entities=False, no_network=True))
            base = posixpath.dirname(posixpath.dirname(name)) if name != '_rels/.rels' else ''
            normalized = False
            for rel in root:
                if rel.get('TargetMode') == 'External':
                    continue
                uri = urlsplit(rel.get('Target', ''))
                resolved = relationship_target(base, rel.get('Target', ''))
                target = '/' + resolved if uri.path.startswith('/') else posixpath.relpath(resolved, base or '.')
                if target != unquote(uri.path):
                    rel.set('Target', urlunsplit(('', '', quote(target, safe='/'), uri.query, uri.fragment)))
                    normalized = True
            if normalized:
                changed[name] = etree.tostring(root, xml_declaration=True, encoding='UTF-8')
        if not changed:
            return source, 0
        directory = out / 'render-input'
        directory.mkdir(exist_ok=True)
        target = directory / source.name
        with zipfile.ZipFile(target, 'w') as copy:
            for entry in archive.infolist():
                copy.writestr(entry, changed.get(entry.filename, archive.read(entry.filename)))
    return target, len(changed)


def validate(source, out):
    started = time.monotonic()
    out.mkdir(parents=True, exist_ok=True)
    report = {'source_sha256': hashlib.sha256(source.read_bytes()).hexdigest(), 'checks': [], 'page_count': 0, 'renderer': 'LibreOffice', 'renderer_version': '', 'pdf_path': None}
    def check(name, action):
        try:
            detail = action()
            report['checks'].append({'name': name, 'status': 'passed', 'issues': [], 'detail': detail})
            return True
        except Exception as error:
            report['checks'].append({'name': name, 'status': 'unavailable' if isinstance(error, (FileNotFoundError, ImportError)) else 'failed', 'issues': [str(error)[:4000]]})
            return False
    package_ok = check('ooxml_package', lambda: {'parts': package_checks(source)})
    if package_ok:
        if source.suffix.lower() == '.docx':
            check('docx_schema', lambda: run(['docx', 'validate', str(source), '--json'])[-4000:])
        elif source.suffix.lower() == '.pptx':
            def read_deck():
                from pptx import Presentation
                deck = Presentation(str(source))
                if not len(deck.slides):
                    raise ValueError('Presentation has no slides')
                return {'slides': len(deck.slides)}
            check('pptx_readability', read_deck)
        else:
            def read_book():
                from openpyxl import load_workbook
                book = load_workbook(source, read_only=True, data_only=False)
                try:
                    return {'sheets': book.sheetnames}
                finally:
                    book.close()
            check('xlsx_readability', read_book)
        profile = out / 'profile'
        conversion_input = None
        def input_for_conversion():
            nonlocal conversion_input
            if conversion_input is None:
                conversion_input = conversion_copy(source, out)
            return conversion_input
        def render():
            report['renderer_version'] = run(['libreoffice', '--version']).strip()
            render_source, normalized_parts = input_for_conversion()
            pdf = convert(render_source, out, profile, 'pdf')
            info = run(['pdfinfo', str(pdf)])
            match = re.search(r'^Pages:\s+(\d+)', info, re.M)
            if not match or int(match[1]) < 1:
                raise ValueError('Rendered document has no pages')
            report['page_count'] = int(match[1])
            if source.suffix.lower() == '.pptx':
                from pptx import Presentation
                if report['page_count'] != len(Presentation(str(source)).slides):
                    raise ValueError('Rendered slide count differs from the saved deck')
            report['pdf_path'] = str(pdf)
            return {'pages': report['page_count'], 'normalized_relationship_parts': normalized_parts}
        check('actual_file_render', render)
        if source.suffix.lower() == '.xlsx':
            def formulas():
                from openpyxl import load_workbook
                recalc = out / 'recalculated'
                recalc.mkdir(exist_ok=True)
                target = convert(input_for_conversion()[0], recalc, profile, 'xlsx:Calc MS Excel 2007 XML')
                original = load_workbook(source, read_only=True, data_only=False)
                calculated = load_workbook(target, read_only=True, data_only=True)
                count = 0
                try:
                    for sheet in original:
                        for row in sheet:
                            for cell in row:
                                if cell.data_type == 'f':
                                    count += 1
                                    result = calculated[sheet.title][cell.coordinate]
                                    if result.data_type == 'e':
                                        raise ValueError(f'Formula failed: {sheet.title}!{cell.coordinate}: {result.value}')
                    return {'formulas_checked': count, 'engine': 'LibreOffice Calc'}
                finally:
                    original.close()
                    calculated.close()
            check('xlsx_formulas', formulas)
        shutil.rmtree(profile, ignore_errors=True)
    else:
        dependent = {'docx': ['docx_schema'], 'pptx': ['pptx_readability'], 'xlsx': ['xlsx_readability', 'xlsx_formulas']}[source.suffix.lower()[1:]]
        for name in [*dependent, 'actual_file_render']:
            report['checks'].append({'name': name, 'status': 'not_run', 'issues': ['Package integrity must pass first']})
    report['elapsed_ms'] = round((time.monotonic() - started) * 1000)
    report['peak_child_rss_kib'] = resource.getrusage(resource.RUSAGE_CHILDREN).ru_maxrss
    print(json.dumps(report))


def pages(pdf, out, selected):
    out.mkdir(parents=True, exist_ok=True)
    for page in selected:
        run(['pdftoppm', '-f', str(page), '-l', str(page), '-singlefile', '-scale-to', '2000', '-jpeg', '-jpegopt', 'quality=90', str(pdf), str(out / f'page-{page}')])
    print(json.dumps({'pages': selected}))


if __name__ == '__main__':
    parser = argparse.ArgumentParser()
    parser.add_argument('action', choices=['validate', 'pages'])
    parser.add_argument('source', type=Path)
    parser.add_argument('output', type=Path)
    parser.add_argument('pages', type=int, nargs='*')
    args = parser.parse_args()
    try:
        if args.action == 'validate':
            with open(f'/tmp/ai-tg-office-render-{os.getuid()}.lock', 'w') as lock:
                fcntl.flock(lock, fcntl.LOCK_EX)
                validate(args.source.resolve(), args.output.resolve())
        else:
            pages(args.source.resolve(), args.output.resolve(), args.pages)
    except Exception as error:
        print(json.dumps({'error': str(error)}))
        sys.exit(1)
