"""Package and actual-render regressions. Runs inside the pinned Office contract."""
import contextlib
import io
import json
from pathlib import Path
import runpy
import sys
import threading
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from unittest.mock import patch
from zipfile import ZipFile, ZIP_DEFLATED

from lxml import etree
from pptx import Presentation

office = runpy.run_path(sys.argv[2] if len(sys.argv) > 2 else '/usr/local/bin/office-files')
root = Path(sys.argv[1]).resolve()
relationship_ns = 'http://schemas.openxmlformats.org/package/2006/relationships'
office_ns = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/'


def rewrite(source, destination, change):
    with ZipFile(source) as archive:
        parts = {name: archive.read(name) for name in archive.namelist()}
    change(parts)
    with ZipFile(destination, 'w', ZIP_DEFLATED) as archive:
        for name, data in parts.items():
            archive.writestr(name, data)
    return destination


def report(source, output):
    captured = io.StringIO()
    with contextlib.redirect_stdout(captured):
        office['validate'](source, output)
    return json.loads(captured.getvalue())


for extension, main in [('docx', 'word/document.xml'), ('pptx', 'ppt/presentation.xml'), ('xlsx', 'xl/workbook.xml')]:
    for prefix in ['/', './unused/../', './']:
        def normalize(parts):
            rels = etree.fromstring(parts['_rels/.rels'])
            for rel in rels:
                if rel.get('Type', '').endswith('/officeDocument'):
                    rel.set('Target', prefix + main)
            parts['_rels/.rels'] = etree.tostring(rels)
        variant = rewrite(root / f'contract.{extension}', root / f'normalized.{extension}', normalize)
        office['package_checks'](variant)

# Pass the original noncanonical targets through all real readers and converters.
for extension in ['docx', 'pptx', 'xlsx']:
    source = root / f'normalized.{extension}'
    before = source.read_bytes()
    output = root / f'normalized-{extension}-qa'
    normalized = report(source, output)
    assert all(check['status'] == 'passed' for check in normalized['checks']), normalized
    assert source.read_bytes() == before
    with ZipFile(source) as original, ZipFile(output / 'render-input' / source.name) as prepared:
        assert original.namelist() == prepared.namelist()
        for name in original.namelist():
            if name != '_rels/.rels':
                assert original.read(name) == prepared.read(name), name

requests = []


class Handler(BaseHTTPRequestHandler):
    def do_GET(self):
        requests.append(self.path)
        self.send_response(200)
        self.end_headers()
        self.wfile.write(b'external resource must not be fetched')

    def log_message(self, *args):
        pass


server = ThreadingHTTPServer(('127.0.0.1', 0), Handler)
thread = threading.Thread(target=server.serve_forever, daemon=True)
thread.start()
url = f'http://127.0.0.1:{server.server_port}/resource'
try:
    for kind in ['image', 'attachedTemplate', 'externalLinkPath', 'unknown', 'hyperlink']:
        def external(parts):
            rels = etree.fromstring(parts['_rels/.rels'])
            etree.SubElement(rels, '{' + relationship_ns + '}Relationship', {
                'Id': 'external-resource', 'Type': office_ns + kind,
                'TargetMode': 'External', 'Target': 'file:///etc/passwd' if kind == 'hyperlink' else url,
            })
            parts['_rels/.rels'] = etree.tostring(rels)
        variant = rewrite(root / 'contract.pptx', root / f'external-{kind}.pptx', external)
        # A rejected package must never reach any converter, even if the renderer
        # is configured to update links or changes its behavior in a later release.
        with patch.dict(office['validate'].__globals__, convert=lambda *args: (_ for _ in ()).throw(AssertionError('conversion invoked'))):
            result = report(variant, root / f'external-{kind}-qa')
        checks = {check['name']: check for check in result['checks']}
        assert checks['ooxml_package']['status'] == 'failed', result
        assert 'External resource' in checks['ooxml_package']['issues'][0], result
        assert checks['actual_file_render']['status'] == 'not_run', result

    # Calling a linked image relationship a hyperlink must not evade the policy.
    def disguised_image(parts):
        name = 'ppt/slides/_rels/slide1.xml.rels'
        rels = etree.fromstring(parts[name])
        rel = next(rel for rel in rels if rel.get('Type', '').endswith('/image'))
        rel.set('Type', office_ns + 'hyperlink')
        rel.set('TargetMode', 'External')
        rel.set('Target', url)
        parts[name] = etree.tostring(rels)
    disguised = rewrite(root / 'contract.pptx', root / 'disguised-image.pptx', disguised_image)
    result = report(disguised, root / 'disguised-qa')
    assert result['checks'][0]['status'] == 'failed', result
    assert 'outside a hyperlink' in result['checks'][0]['issues'][0], result

    # Preserve ordinary clickable links and prove rendering them makes no request.
    presentation = Presentation(str(root / 'contract.pptx'))
    run = presentation.slides[0].shapes[0].text_frame.paragraphs[0].runs[0]
    run.hyperlink.address = url
    linked = root / 'hyperlink.pptx'
    presentation.save(str(linked))
    result = report(linked, root / 'hyperlink-qa')
    assert all(check['status'] == 'passed' for check in result['checks']), result
    assert not requests, requests
finally:
    server.shutdown()
    server.server_close()
    thread.join()

print('Normalized roots, blank formulas, blocked external resources, and inert hyperlinks passed')
